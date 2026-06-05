# -*- coding: utf-8 -*-
"""
사내 연차 경매 시스템 — 발표용 빈 PPTX 템플릿 생성기 (v3: 라이트 / 하늘색 톤)
============================================================================
디자인 방향 (제품 로그인 화면과 동일한 결)
  - 밝고 가벼운 톤: 화이트 ~ 옅은 하늘색 배경, 부드러운 원으로 깊이감.
  - pill 배지(B2E · ESCROW & DIVIDEND), 화이트 지표 카드(라벨/수치/설명).
  - 큰 제목 + 짧은 파랑 액센트 라인, 좌상단 로고, 하단 헤어라인 푸터.
  - 전 요소 네이티브(편집 가능). placeholder 는 lstStyle 상속으로 스타일링.

팔레트 (frontend globals.css 메인색 + 라이트 톤 보강)
  MAIN #1B67DA · DEEP #0E4EAF · 페이지배경 #F4F8FF · 카드 #FFFFFF · pill #E1ECFF

실행:  python presentation/template/build-template.py
산출물: 타임소프트콘_템플릿.pptx
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree

# ---- 팔레트 -----------------------------------------------------------------
MAIN   = "1B67DA"   # 프로젝트 메인 (--primary)
DEEP   = "0E4EAF"   # 딥 블루 (카드 수치/로고 등 액센트)
TITLE  = "14315F"   # 제목용 딥네이비 (차분·고급)
MID    = "4F86E8"   # 중간 파랑
PAGEBG = "F4F8FF"   # 페이지 배경(옅은 하늘색)
PAGEBG2 = "E8F0FF"  # 그라데이션 끝
CONTENTBG = "EEF4FF"  # 일반 본문 슬라이드 배경(하늘색 틴트)
CIRCLE = "DCE8FF"   # 배경 장식 원
CARDBG = "FFFFFF"
LINE   = "E5EBF4"   # 카드/구분 라인
PILLBG = "E1ECFF"   # 배지 배경
PILLTX = "1B57C0"   # 배지 글자
INK    = "1B2433"   # 제목/본문 진한 글자
BODY   = "39414E"   # 본문
GRAY   = "6B7280"   # 보조
LABEL  = "94A0B3"   # 카드 라벨/푸터
HAIR   = "D8E2F2"   # 헤어라인(틴트 배경에서도 보이게)

KO_FONT = "Malgun Gothic"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"

T_TITLE, T_BODY, T_CTITLE, T_SUBT, T_OBJ, T_PIC = 1, 2, 3, 4, 7, 18
CONTENT_TYPES = (T_BODY, T_OBJ, T_PIC)
MARGIN = Inches(0.9)
_id = [1000]


def sub(parent, tag, **attrs):
    el = etree.SubElement(parent, qn(tag))
    for k, v in attrs.items():
        el.set(k, v)
    return el


def _sptree(s):
    return s.element.find(qn("p:cSld")).find(qn("p:spTree"))


def _nid():
    _id[0] += 1
    return _id[0]


# ============================================================ 도형 헬퍼
def add_rect(s, x, y, w, h, hexv, alpha=None, rounded=False, radius=25000,
             line_hex=None, line_w=None, shadow=False, behind=False):
    tree = _sptree(s)
    sp = etree.SubElement(tree, qn("p:sp"))
    nv = sub(sp, "p:nvSpPr")
    c = sub(nv, "p:cNvPr"); c.set("id", str(_nid())); c.set("name", "shp")
    sub(nv, "p:cNvSpPr"); sub(nv, "p:nvPr")
    pr = sub(sp, "p:spPr")
    xf = sub(pr, "a:xfrm")
    sub(xf, "a:off", x=str(int(x)), y=str(int(y)))
    sub(xf, "a:ext", cx=str(int(w)), cy=str(int(h)))
    g = sub(pr, "a:prstGeom", prst="roundRect" if rounded else "rect")
    av = sub(g, "a:avLst")
    if rounded:
        sub(av, "a:gd", name="adj", fmla="val %d" % radius)
    fill = sub(pr, "a:solidFill")
    col = sub(fill, "a:srgbClr", val=hexv)
    if alpha is not None:
        sub(col, "a:alpha", val=str(int(alpha * 1000)))
    if line_hex:
        ln = sub(pr, "a:ln", w=str(int(line_w or Pt(1))))
        lf = sub(ln, "a:solidFill"); sub(lf, "a:srgbClr", val=line_hex)
    else:
        ln = sub(pr, "a:ln"); sub(ln, "a:noFill")
    if shadow:
        eff = sub(pr, "a:effectLst")
        sh = sub(eff, "a:outerShdw", blurRad="90000", dist="38100",
                 dir="5400000", rotWithShape="0")
        sc = sub(sh, "a:srgbClr", val="1B3A6B"); sub(sc, "a:alpha", val="14000")
    tx = sub(sp, "p:txBody"); sub(tx, "a:bodyPr"); sub(tx, "a:lstStyle"); sub(tx, "a:p")
    if behind:
        # spTree 앞쪽(nvGrpSpPr, grpSpPr 다음)으로 옮겨 placeholder 뒤에 깔리게
        tree.remove(sp); tree.insert(2, sp)
    return sp


def set_bg_solid(s, hexv):
    cSld = s.element.find(qn("p:cSld"))
    old = cSld.find(qn("p:bg"))
    if old is not None:
        cSld.remove(old)
    bg = etree.Element(qn("p:bg")); bgpr = sub(bg, "p:bgPr")
    f = sub(bgpr, "a:solidFill"); sub(f, "a:srgbClr", val=hexv)
    sub(bgpr, "a:effectLst")
    cSld.insert(0, bg)


def add_oval(s, x, y, w, h, fill_hex=None, fill_alpha=None,
             line_hex=None, line_w=None, line_alpha=None):
    sp = etree.SubElement(_sptree(s), qn("p:sp"))
    nv = sub(sp, "p:nvSpPr")
    c = sub(nv, "p:cNvPr"); c.set("id", str(_nid())); c.set("name", "circle")
    sub(nv, "p:cNvSpPr"); sub(nv, "p:nvPr")
    pr = sub(sp, "p:spPr")
    xf = sub(pr, "a:xfrm")
    sub(xf, "a:off", x=str(int(x)), y=str(int(y)))
    sub(xf, "a:ext", cx=str(int(w)), cy=str(int(h)))
    g = sub(pr, "a:prstGeom", prst="ellipse"); sub(g, "a:avLst")
    if fill_hex:
        f = sub(pr, "a:solidFill"); col = sub(f, "a:srgbClr", val=fill_hex)
        if fill_alpha is not None:
            sub(col, "a:alpha", val=str(int(fill_alpha * 1000)))
    else:
        sub(pr, "a:noFill")
    if line_hex:
        ln = sub(pr, "a:ln", w=str(int(line_w))) if line_w else sub(pr, "a:ln")
        f = sub(ln, "a:solidFill"); col = sub(f, "a:srgbClr", val=line_hex)
        if line_alpha is not None:
            sub(col, "a:alpha", val=str(int(line_alpha * 1000)))
    else:
        ln = sub(pr, "a:ln"); sub(ln, "a:noFill")
    tx = sub(sp, "p:txBody"); sub(tx, "a:bodyPr"); sub(tx, "a:lstStyle"); sub(tx, "a:p")
    return sp


def add_text(s, x, y, w, h, text, size, color, bold=False,
             align="l", spc=None, anchor="ctr", caps=False, line_spc=None):
    sp = etree.SubElement(_sptree(s), qn("p:sp"))
    nv = sub(sp, "p:nvSpPr")
    c = sub(nv, "p:cNvPr"); c.set("id", str(_nid())); c.set("name", "txt")
    sub(nv, "p:cNvSpPr"); sub(nv, "p:nvPr")
    pr = sub(sp, "p:spPr")
    xf = sub(pr, "a:xfrm")
    sub(xf, "a:off", x=str(int(x)), y=str(int(y)))
    sub(xf, "a:ext", cx=str(int(w)), cy=str(int(h)))
    g = sub(pr, "a:prstGeom", prst="rect"); sub(g, "a:avLst")
    sub(pr, "a:noFill"); ln = sub(pr, "a:ln"); sub(ln, "a:noFill")
    tx = sub(sp, "p:txBody")
    bp = sub(tx, "a:bodyPr", wrap="square", anchor=anchor)
    bp.set("lIns", "0"); bp.set("rIns", "0")
    sub(tx, "a:lstStyle")
    for li, line in enumerate(text.split("\n")):
        para = sub(tx, "a:p")
        ppr = sub(para, "a:pPr", algn=align)
        if line_spc is not None:
            ls = sub(ppr, "a:lnSpc"); sub(ls, "a:spcPct", val=str(int(line_spc * 1000)))
        r = sub(para, "a:r")
        rpr = sub(r, "a:rPr", lang="ko-KR", sz=str(int(size * 100)))
        rpr.set("b", "1" if bold else "0")
        if spc is not None:
            rpr.set("spc", str(int(spc * 100)))
        f = sub(rpr, "a:solidFill"); sub(f, "a:srgbClr", val=color)
        sub(rpr, "a:latin", typeface=KO_FONT); sub(rpr, "a:ea", typeface=KO_FONT)
        t = sub(r, "a:t"); t.text = (line.upper() if caps else line)
    return sp


def add_logo(s, x, y, scale=1.0):
    """좌상단 로고 마크(둥근 사각 + 점) + 워드마크(라이트 배경용 진한 글자)."""
    sz = Inches(0.24 * scale)
    add_rect(s, x, y, sz, sz, MAIN, rounded=True, radius=30000)
    add_rect(s, x + sz * 0.30, y + sz * 0.55, sz * 0.42, sz * 0.42,
             "FFFFFF", rounded=True, radius=40000)
    add_text(s, x + sz + Inches(0.14), y - Inches(0.02), Inches(4),
             sz + Inches(0.04), "타임소프트콘", 13 * scale, DEEP, bold=True,
             anchor="ctr", spc=0.5)


def add_pill(s, x, y, text):
    w = Inches(0.42) + Inches(0.092) * len(text)
    h = Inches(0.36)
    add_rect(s, x, y, w, h, PILLBG, rounded=True, radius=50000)
    add_text(s, x, y, w, h, text, 11, PILLTX, bold=True,
             align="ctr", anchor="ctr", spc=1.5, caps=True)
    return w


def add_card(s, x, y, w, h, label, number, subtxt):
    add_rect(s, x, y, w, h, CARDBG, rounded=True, radius=14000,
             line_hex=LINE, line_w=Pt(1), shadow=True)
    pad = Inches(0.24)
    add_text(s, x + pad, y + Inches(0.16), w - 2 * pad, Inches(0.26),
             label, 10.5, LABEL, bold=True, anchor="t", spc=0.5)
    add_text(s, x + pad, y + Inches(0.40), w - 2 * pad, Inches(0.52),
             number, 25, DEEP, bold=True, anchor="t")
    add_text(s, x + pad, y + h - Inches(0.36), w - 2 * pad, Inches(0.26),
             subtxt, 9.5, GRAY, anchor="t")


# ============================================================ 배경
def set_bg_gradient(s, stops, ang_deg):
    cSld = s.element.find(qn("p:cSld"))
    old = cSld.find(qn("p:bg"))
    if old is not None:
        cSld.remove(old)
    bg = etree.Element(qn("p:bg")); bgpr = sub(bg, "p:bgPr")
    grad = sub(bgpr, "a:gradFill"); gs = sub(grad, "a:gsLst")
    for pos, hexv in stops:
        g = sub(gs, "a:gs", pos=str(int(pos * 1000)))
        sub(g, "a:srgbClr", val=hexv)
    sub(grad, "a:lin", ang=str(int(ang_deg * 60000)), scaled="1")
    sub(bgpr, "a:effectLst")
    cSld.insert(0, bg)


# ============================================================ placeholder 스타일
def style_ph(ph, size=None, color=None, bold=None, align=None, spc=None,
             font=KO_FONT, line_spc=None, bullet=True):
    txb = ph.text_frame._txBody
    lst = txb.find(qn("a:lstStyle"))
    if lst is None:
        lst = etree.Element(qn("a:lstStyle"))
        txb.find(qn("a:bodyPr")).addnext(lst)
    for e in lst.findall(qn("a:lvl1pPr")):
        lst.remove(e)
    lvl = etree.Element(qn("a:lvl1pPr"))
    lst.insert(0, lvl)
    if align:
        lvl.set("algn", align)
    if line_spc is not None:
        ls = sub(lvl, "a:lnSpc"); sub(ls, "a:spcPct", val=str(int(line_spc * 1000)))
    if not bullet:
        lvl.set("indent", "0"); lvl.set("marL", "0"); sub(lvl, "a:buNone")
    defr = sub(lvl, "a:defRPr")
    if size is not None:
        defr.set("sz", str(int(size * 100)))
    if bold is not None:
        defr.set("b", "1" if bold else "0")
    if spc is not None:
        defr.set("spc", str(int(spc * 100)))
    if color is not None:
        f = sub(defr, "a:solidFill"); sub(f, "a:srgbClr", val=color)
    sub(defr, "a:latin", typeface=font); sub(defr, "a:ea", typeface=font)


def move_ph(ph, x, y, w, h):
    ph.left, ph.top, ph.width, ph.height = int(x), int(y), int(w), int(h)


# ============================================================ theme / master txstyles
def restyle_theme(prs):
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    master = prs.slide_masters[0]
    tp = master.part.part_related_by(RT.THEME)
    root = etree.fromstring(tp.blob)
    te = root.find(qn("a:themeElements"))
    old = te.find(qn("a:clrScheme")); te.remove(old)
    clr = etree.Element(qn("a:clrScheme")); clr.set("name", "타임소프트콘")

    def srgb(tag, hexv):
        w = sub(clr, tag); sub(w, "a:srgbClr", val=hexv)

    def sysc(tag, val, last):
        w = sub(clr, tag); cc = sub(w, "a:sysClr", val=val); cc.set("lastClr", last)

    sysc("a:dk1", "windowText", "000000"); sysc("a:lt1", "window", "FFFFFF")
    srgb("a:dk2", DEEP); srgb("a:lt2", PAGEBG)
    srgb("a:accent1", MAIN); srgb("a:accent2", DEEP); srgb("a:accent3", MID)
    srgb("a:accent4", "11366E"); srgb("a:accent5", "2BB7A3"); srgb("a:accent6", GRAY)
    srgb("a:hlink", MAIN); srgb("a:folHlink", "11366E")
    te.insert(0, clr)

    fs = te.find(qn("a:fontScheme")); fs.set("name", "타임소프트콘")
    for grp in ("a:majorFont", "a:minorFont"):
        g = fs.find(qn(grp))
        g.find(qn("a:latin")).set("typeface", KO_FONT)
        ea = g.find(qn("a:ea"))
        (ea if ea is not None else sub(g, "a:ea")).set("typeface", KO_FONT)
    tp._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)


def restyle_master_txstyles(prs):
    master = prs.slide_masters[0]
    txs = master.element.find(qn("p:txStyles"))
    BODY_SZ = [18, 16, 15, 14, 13, 13, 13, 13, 13]

    def paint(dr, color, bold=None, size=None):
        if bold is not None:
            dr.set("b", "1" if bold else "0")
        if size is not None:
            dr.set("sz", str(int(size * 100)))
        for f in dr.findall(qn("a:solidFill")):
            dr.remove(f)
        fill = etree.Element(qn("a:solidFill")); sub(fill, "a:srgbClr", val=color)
        dr.insert(0, fill)
        for ft in ("a:latin", "a:ea"):
            el = dr.find(qn(ft))
            if el is None:
                el = sub(dr, ft)
            el.set("typeface", KO_FONT)

    ts = txs.find(qn("p:titleStyle"))
    if ts is not None:
        for lvl in ts:
            dr = lvl.find(qn("a:defRPr"))
            if dr is not None:
                paint(dr, TITLE, bold=True); break
    bs = txs.find(qn("p:bodyStyle"))
    if bs is not None:
        for i, lvl in enumerate(bs):
            dr = lvl.find(qn("a:defRPr"))
            if dr is not None:
                paint(dr, BODY, size=BODY_SZ[min(i, len(BODY_SZ) - 1)])


def decorate_master(prs):
    m = prs.slide_masters[0]
    sw, sh = prs.slide_width, prs.slide_height
    set_bg_solid(m, CONTENTBG)   # 일반 본문 = 하늘색 틴트 (표지/섹션은 레이아웃에서 오버라이드)
    add_logo(m, MARGIN, Inches(0.42), scale=0.92)
    add_rect(m, MARGIN, sh - Inches(0.58), sw - 2 * MARGIN, Pt(1), HAIR)
    add_text(m, MARGIN, sh - Inches(0.52), Inches(7), Inches(0.3),
             "사내 연차 경매 시스템", 9, LABEL, anchor="ctr", spc=0.3)


# ============================================================ 레이아웃
LAYOUT_NAMES = {
    0: "표지", 1: "제목+본문", 2: "섹션 구분", 3: "2단 본문",
    4: "비교", 5: "제목만", 6: "빈 화면", 7: "본문+캡션",
    8: "이미지+캡션", 9: "제목+세로글", 10: "세로 제목+글",
}


def hide_master(layout):
    layout.element.set("showMasterSp", "0")


def title_block(layout, sw):
    for ph in layout.placeholders:
        if ph.placeholder_format.type in (T_TITLE, T_CTITLE):
            move_ph(ph, MARGIN, Inches(0.62), sw - 2 * MARGIN, Inches(0.95))
            style_ph(ph, size=29, color=TITLE, bold=True, align="l")
            add_rect(layout, MARGIN, Inches(1.52), Inches(0.62), Pt(3.5), MAIN)
            return ph
    return None


def content_phs(layout):
    return [ph for ph in layout.placeholders
            if ph.placeholder_format.type in CONTENT_TYPES]


def soft_circles(s, sw, sh):
    """우측에 부드러운 옅은 원 2개(제품 화면 결)."""
    add_oval(s, sw - Inches(4.2), Inches(0.6), Inches(5.0), Inches(5.0),
             fill_hex=CIRCLE, fill_alpha=55)
    add_oval(s, sw - Inches(2.2), Inches(2.7), Inches(3.4), Inches(3.4),
             fill_hex=MID, fill_alpha=10)


def restyle_layouts(prs):
    sw, sh = prs.slide_width, prs.slide_height
    full_w = sw - 2 * MARGIN
    c_top, c_h = Inches(1.85), Inches(4.9)

    for idx, layout in enumerate(prs.slide_layouts):
        if idx in LAYOUT_NAMES:
            layout.name = LAYOUT_NAMES[idx]

        # ---------- 표지 (라이트) ----------
        if idx == 0:
            hide_master(layout)
            set_bg_gradient(layout, [(0, "FFFFFF"), (100, PAGEBG2)], 65)
            soft_circles(layout, sw, sh)
            add_logo(layout, MARGIN, Inches(0.6), scale=1.1)
            add_rect(layout, MARGIN, Inches(3.95), Inches(0.9), Pt(4), MAIN)
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                if t == T_CTITLE:
                    move_ph(ph, MARGIN - Inches(0.02), Inches(2.55),
                            Inches(9.5), Inches(1.4))
                    style_ph(ph, size=46, color=INK, bold=True, align="l")
                elif t == T_SUBT:
                    move_ph(ph, MARGIN, Inches(4.2), Inches(9), Inches(0.9))
                    style_ph(ph, size=18, color=GRAY, bold=False, align="l",
                             line_spc=132, bullet=False)
            continue

        # ---------- 섹션 구분 (라이트) ----------
        if idx == 2:
            hide_master(layout)
            set_bg_gradient(layout, [(0, "FFFFFF"), (100, PAGEBG2)], 65)
            soft_circles(layout, sw, sh)
            add_logo(layout, MARGIN, Inches(0.55), scale=0.95)
            add_text(layout, MARGIN, Inches(2.45), Inches(6), Inches(0.4),
                     "SECTION", 13, MAIN, bold=True, spc=3.0, caps=True)
            add_rect(layout, MARGIN, Inches(4.35), Inches(0.9), Pt(4), MAIN)
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                if t == T_TITLE:
                    move_ph(ph, MARGIN, Inches(2.9), Inches(10), Inches(1.4))
                    style_ph(ph, size=40, color=TITLE, bold=True, align="l")
                elif t == T_BODY:
                    move_ph(ph, MARGIN, Inches(4.65), Inches(9), Inches(0.9))
                    style_ph(ph, size=18, color=GRAY, bold=False, align="l",
                             bullet=False)
            continue

        if idx == 6:   # 빈 화면
            continue

        # ---------- 이미지+캡션: 흰 카드 위에 그림 ----------
        if idx == 8:
            title_block(layout, sw)
            pics = [p for p in layout.placeholders if p.placeholder_format.type == T_PIC]
            caps = [p for p in layout.placeholders if p.placeholder_format.type == T_BODY]
            cx, cy, cw, ch = MARGIN, Inches(1.95), Inches(7.55), Inches(4.75)
            inset = Inches(0.16)
            add_rect(layout, cx, cy, cw, ch, CARDBG, rounded=True, radius=9000,
                     line_hex=LINE, line_w=Pt(1), shadow=True, behind=True)
            if pics:
                move_ph(pics[0], cx + inset, cy + inset,
                        cw - 2 * inset, ch - 2 * inset)
            if caps:
                rx = cx + cw + Inches(0.5)
                move_ph(caps[0], rx, Inches(1.95), sw - MARGIN - rx, Inches(4.75))
                style_ph(caps[0], size=16, color=BODY, align="l",
                         line_spc=132, bullet=False)
            continue

        # ---------- 일반 레이아웃 ----------
        title_block(layout, sw)
        cps = content_phs(layout)
        if len(cps) == 1:
            ph = cps[0]
            move_ph(ph, MARGIN, c_top, full_w, c_h)
            if ph.placeholder_format.type != T_PIC:
                style_ph(ph, size=18, color=BODY, align="l", line_spc=128)
        elif len(cps) == 2:
            gap = Inches(0.5); half = (full_w - gap) / 2
            for i, ph in enumerate(cps):
                x = MARGIN + (0 if i == 0 else half + gap)
                move_ph(ph, x, c_top, half, c_h)
                if ph.placeholder_format.type != T_PIC:
                    style_ph(ph, size=17, color=BODY, align="l", line_spc=126)
        elif len(cps) >= 3:
            gap = Inches(0.5); half = (full_w - gap) / 2
            left = [p for p in cps if p.left is None or p.left < sw / 2]
            right = [p for p in cps if p not in left]
            for col, base_x in ((left, MARGIN), (right, MARGIN + half + gap)):
                y = c_top
                for ph in col:
                    is_hdr = ph.placeholder_format.type == T_BODY
                    hh = Inches(0.55) if is_hdr else c_h - Inches(0.7)
                    move_ph(ph, base_x, y, half, hh)
                    style_ph(ph, size=17 if is_hdr else 16, color=BODY,
                             bold=is_hdr, align="l", line_spc=124)
                    y = y + hh + Inches(0.15)


# ============================================================ 샘플 슬라이드
def add_samples(prs):
    sw, sh = prs.slide_width, prs.slide_height

    # 1) 표지 — pill 배지 + 하단 지표 카드(편집 가능)
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_pill(s, MARGIN, Inches(2.0), "B2E · ESCROW & DIVIDEND")
    for ph in s.placeholders:
        t = ph.placeholder_format.type
        if t == T_CTITLE:
            ph.text = "타임소프트콘"
        elif t == T_SUBT:
            ph.text = "사내 연차 경매 시스템 · 김기철 · 오지석"
    cards = [("재무 리스크", "0%", "회사 예산 신규투입 없음"),
             ("재무 정합성", "100%", "Insert-Only 원장"),
             ("동시성 제어", "SQLite", "write 락 (직렬화)")]
    cw = Inches(3.62); gap = Inches(0.33); cy = sh - Inches(1.75); ch = Inches(1.18)
    for i, (lab, num, sb) in enumerate(cards):
        add_card(s, MARGIN + i * (cw + gap), cy, cw, ch, lab, num, sb)

    # 2) 섹션
    s = prs.slides.add_slide(prs.slide_layouts[2])
    for ph in s.placeholders:
        t = ph.placeholder_format.type
        if t == T_TITLE:
            ph.text = "01  문제 정의"
        elif t == T_BODY:
            ph.text = "왜 P2P가 아니라 B2E 에스크로인가"

    # 3) 제목+본문
    s = prs.slides.add_slide(prs.slide_layouts[1])
    for ph in s.placeholders:
        t = ph.placeholder_format.type
        if t == T_TITLE:
            ph.text = "편집 가능한 네이티브 슬라이드"
        elif t in (T_OBJ, T_BODY):
            tf = ph.text_frame
            tf.text = "모든 텍스트는 PowerPoint에서 직접 수정됩니다"
            for txt, lvl in [("테마색·폰트는 슬라이드 마스터에서 일괄 관리", 0),
                             ("새 슬라이드: 홈 ▸ 새 슬라이드 ▸ 레이아웃 선택", 0),
                             ("제목 옆 짧은 라인이 브랜드 액센트", 1)]:
                p = tf.add_paragraph(); p.text = txt; p.level = lvl

    # 4) 2단 본문
    s = prs.slides.add_slide(prs.slide_layouts[3])
    for ph in s.placeholders:
        t = ph.placeholder_format.type
        if t == T_TITLE:
            ph.text = "좌우 2단 레이아웃"
        elif t in (T_OBJ, T_BODY):
            ph.text_frame.text = "왼쪽 또는 오른쪽 본문 영역"

    # 5) 이미지+캡션 — 틴트 배경 위 흰 카드 데모(실제 프로젝트 그림 사용)
    here = os.path.dirname(os.path.abspath(__file__))
    img = None
    for cand in ("architecture-block.png", "clean-arch.png", "subsystems.png"):
        p = os.path.join(here, "..", "img", cand)
        if os.path.exists(p):
            img = p; break
    s = prs.slides.add_slide(prs.slide_layouts[8])
    phs = {ph.placeholder_format.type: ph for ph in s.placeholders}
    if T_TITLE in phs:
        phs[T_TITLE].text = "이미지·스크린샷은 흰 카드 위에"
    if T_BODY in phs:
        phs[T_BODY].text = ("틴트 배경에서도 흰 배경 그림이 둥둥 뜨지 않고\n"
                            "의도된 카드처럼 보입니다.\n\n캡션·설명은 여기에 적습니다.")
    if T_PIC in phs and img:
        phs[T_PIC].insert_picture(img)


# ============================================================ main
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    restyle_theme(prs)
    restyle_master_txstyles(prs)
    decorate_master(prs)
    restyle_layouts(prs)
    add_samples(prs)
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "타임소프트콘_템플릿.pptx")
    prs.save(out)
    print("OK ->", out)


if __name__ == "__main__":
    main()
